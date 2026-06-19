"""Helpers for building CANlab code-map .pptx files in the style of
docs/canlab_template_codemap.pptx.

The template is a single-slide widescreen deck (13.33" x 7.5") with:
  - a centered title and description band along the top
  - a left "Use:" / "Notes:" prose column
  - a free-form data-flow diagram in the center using four shape colors
        files          = FCE5BF (light orange)
        canlab funcs   = DAE3F3 (light blue)
        variables      = E2EFD9 (light green)
        custom scripts = FFFFFF, with thin border
  - down/right arrows connecting nodes
  - a small legend in the lower-left

This module exposes a Slide wrapper with primitives (box, arrow, text,
legend, header) so each per-function spec stays compact and can hand-place
shapes for clarity. Coordinates are in inches.
"""

from __future__ import annotations

import os
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from typing import Iterable, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from lxml import etree

# ---- Color palette (RGB hex from the template) -----------------------------
COLOR_FILE = RGBColor(0xFC, 0xE5, 0xBF)
COLOR_FUNC = RGBColor(0xDA, 0xE3, 0xF3)
COLOR_VAR = RGBColor(0xE2, 0xEF, 0xD9)
COLOR_SCRIPT = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BORDER = RGBColor(0x80, 0x80, 0x80)
COLOR_TEXT = RGBColor(0x26, 0x26, 0x26)
COLOR_TEXT_LIGHT = RGBColor(0x59, 0x59, 0x59)

KIND_TO_COLOR = {
    "file": COLOR_FILE,
    "func": COLOR_FUNC,
    "var": COLOR_VAR,
    "script": COLOR_SCRIPT,
}

# ---- Default sizes (inches) ------------------------------------------------
DEFAULT_BOX_W = 1.85
DEFAULT_BOX_H = 0.42
DEFAULT_ARROW_W = 0.30
DEFAULT_ARROW_H = 0.32
DEFAULT_RIGHT_ARROW_W = 0.74
DEFAULT_RIGHT_ARROW_H = 0.30


@dataclass
class Node:
    """A placed node on the slide. Coordinates are top-left in inches."""

    id: str
    kind: str  # 'file' | 'func' | 'var' | 'script'
    label: str
    x: float
    y: float
    w: float = DEFAULT_BOX_W
    h: float = DEFAULT_BOX_H
    font_name: Optional[str] = None  # default: Calibri for func/file/script, Cambria for var
    font_size: float = 11.0
    bold: bool = False

    @property
    def cx(self) -> float:
        return self.x + self.w / 2

    @property
    def cy(self) -> float:
        return self.y + self.h / 2

    @property
    def bottom(self) -> float:
        return self.y + self.h

    @property
    def right(self) -> float:
        return self.x + self.w


class Slide:
    """Thin wrapper around a python-pptx slide for code-map drawing."""

    def __init__(self, title: str, description: Iterable[str]):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333333)
        self.prs.slide_height = Inches(7.5)
        # Use a blank layout (index 6 in default theme)
        blank_layout = self.prs.slide_layouts[6]
        self.slide = self.prs.slides.add_slide(blank_layout)
        self.nodes: dict[str, Node] = {}
        self._draw_header(title, list(description))

    # -- header / prose -----------------------------------------------------
    def _draw_header(self, title: str, description: List[str]) -> None:
        # Title — auto-shrink for long names so it fits on one line.
        # Calibrated against soffice's Calibri Bold rendering at 9.30" wide.
        n = len(title)
        if n <= 45:
            title_size = 22.0
        elif n <= 52:
            title_size = 19.0
        elif n <= 58:
            title_size = 17.0
        elif n <= 66:
            title_size = 15.0
        else:
            title_size = 13.5
        self._text(
            title,
            x=2.00,
            y=0.15,
            w=9.30,
            h=0.45,
            size=title_size,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        # Description (1-3 lines, centered)
        if description:
            self._text(
                "\n".join(description),
                x=3.50,
                y=0.62,
                w=7.00,
                h=0.55 + 0.18 * max(0, len(description) - 2),
                size=11.0,
                align=PP_ALIGN.CENTER,
            )

    def use(self, bullets: Iterable[str]) -> None:
        self._text("Use:", x=0.24, y=0.58, w=0.60, h=0.25, size=12.0, bold=True)
        self._text(
            "\n\n".join(bullets),
            x=0.25,
            y=0.86,
            w=2.55,
            h=2.65,
            size=10.5,
        )

    def notes(self, bullets: Iterable[str], y: float = 3.67, h: float = 1.50) -> None:
        self._text("Notes:", x=0.26, y=y, w=0.70, h=0.25, size=12.0, bold=True)
        self._text(
            "\n\n".join(bullets),
            x=0.26,
            y=y + 0.28,
            w=2.10,
            h=h,
            size=10.5,
        )

    def object_types(self, x: float, y: float, lines: Iterable[str], width: float = 1.7) -> None:
        self._text(
            "Object types:",
            x=x,
            y=y,
            w=width,
            h=0.22,
            size=10.0,
            bold=True,
        )
        self._text(
            "\n".join(lines),
            x=x,
            y=y + 0.20,
            w=width,
            h=0.25 * max(1, len(list(lines)) if not isinstance(lines, list) else len(lines)),
            size=10.0,
        )

    # Object_types accepts an iterable; cache once by reading list
    def object_types_panel(self, x: float, y: float, lines: List[str], width: float = 1.7) -> None:
        self._text("Object types:", x=x, y=y, w=width, h=0.22, size=10.0, bold=True)
        self._text(
            "\n".join(lines),
            x=x,
            y=y + 0.22,
            w=width,
            h=0.22 * len(lines) + 0.06,
            size=10.0,
        )

    # -- nodes / arrows -----------------------------------------------------
    def box(
        self,
        node_id: str,
        kind: str,
        label: str,
        x: float,
        y: float,
        w: float = DEFAULT_BOX_W,
        h: float = DEFAULT_BOX_H,
        font_name: Optional[str] = None,
        font_size: float = 11.0,
        bold: bool = False,
    ) -> Node:
        if kind not in KIND_TO_COLOR:
            raise ValueError(f"Unknown box kind: {kind}")
        node = Node(node_id, kind, label, x, y, w, h, font_name, font_size, bold)
        self.nodes[node_id] = node
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = KIND_TO_COLOR[kind]
        # Border: light gray for everything; visible for white "script" boxes
        line = shape.line
        line.color.rgb = COLOR_BORDER
        line.width = Pt(0.5)
        # Text
        tf = shape.text_frame
        tf.margin_left = Inches(0.04)
        tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = label
        run.font.name = font_name or ("Cambria" if kind == "var" else "Calibri")
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = COLOR_TEXT
        return node

    def arrow_down(
        self,
        x: float,
        y: float,
        h: float = DEFAULT_ARROW_H,
        w: float = DEFAULT_ARROW_W,
    ) -> None:
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_SCRIPT
        shape.line.color.rgb = COLOR_BORDER
        shape.line.width = Pt(0.5)

    def arrow_right(
        self,
        x: float,
        y: float,
        w: float = DEFAULT_RIGHT_ARROW_W,
        h: float = DEFAULT_RIGHT_ARROW_H,
    ) -> None:
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_SCRIPT
        shape.line.color.rgb = COLOR_BORDER
        shape.line.width = Pt(0.5)

    def arrow_left(
        self,
        x: float,
        y: float,
        w: float = DEFAULT_RIGHT_ARROW_W,
        h: float = DEFAULT_RIGHT_ARROW_H,
    ) -> None:
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.LEFT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_SCRIPT
        shape.line.color.rgb = COLOR_BORDER
        shape.line.width = Pt(0.5)

    def arrow_up(
        self,
        x: float,
        y: float,
        h: float = DEFAULT_ARROW_H,
        w: float = DEFAULT_ARROW_W,
    ) -> None:
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.UP_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_SCRIPT
        shape.line.color.rgb = COLOR_BORDER
        shape.line.width = Pt(0.5)

    def connect_down(self, src: str, dst: str, gap: float = 0.05) -> None:
        """Place a small down-arrow centered between two stacked nodes."""
        a = self.nodes[src]
        b = self.nodes[dst]
        x_center = (a.cx + b.cx) / 2 - DEFAULT_ARROW_W / 2
        y_top = a.bottom + gap
        y_bottom = b.y - gap
        self.arrow_down(x_center, y_top, h=max(0.18, y_bottom - y_top))

    def connect_right(self, src: str, dst: str, gap: float = 0.05) -> None:
        a = self.nodes[src]
        b = self.nodes[dst]
        y_center = (a.cy + b.cy) / 2 - DEFAULT_RIGHT_ARROW_H / 2
        x_left = a.right + gap
        x_right = b.x - gap
        self.arrow_right(x_left, y_center, w=max(0.30, x_right - x_left))

    def connect_left(self, src: str, dst: str, gap: float = 0.05) -> None:
        a = self.nodes[src]
        b = self.nodes[dst]
        y_center = (a.cy + b.cy) / 2 - DEFAULT_RIGHT_ARROW_H / 2
        x_right = a.x - gap
        x_left = b.right + gap
        self.arrow_left(x_left, y_center, w=max(0.30, x_right - x_left))

    # -- thin-line connector arrows (for diagonals & elbows) -----------------
    def line_arrow(
        self,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        weight: float = 1.25,
        color: RGBColor = COLOR_BORDER,
        head: str = "end",
        kind: str = "straight",
    ) -> None:
        """Thin connector line with an arrowhead.

        kind='straight' draws a single segment from (x1, y1) → (x2, y2).
        kind='bent' draws an L-shape: horizontal from x1 then vertical to (x2, y2).
        head='end' puts the arrowhead at the destination; 'start' at the origin;
        'both' on both ends.
        """
        if kind == "straight":
            connector_type = MSO_CONNECTOR.STRAIGHT
        elif kind == "bent":
            connector_type = MSO_CONNECTOR.ELBOW
        elif kind == "curve":
            connector_type = MSO_CONNECTOR.CURVE
        else:
            raise ValueError(f"Unknown connector kind: {kind}")

        conn = self.slide.shapes.add_connector(
            connector_type,
            Inches(x1),
            Inches(y1),
            Inches(x2),
            Inches(y2),
        )
        line = conn.line
        line.color.rgb = color
        line.width = Pt(weight)
        # Add arrowheads via XML (python-pptx doesn't expose them directly).
        ln = conn.line._get_or_add_ln()
        if head in ("end", "both"):
            tail_end = ln.find(qn("a:tailEnd"))
            if tail_end is None:
                tail_end = etree.SubElement(ln, qn("a:tailEnd"))
            tail_end.set("type", "triangle")
            tail_end.set("w", "med")
            tail_end.set("len", "med")
        if head in ("start", "both"):
            head_end = ln.find(qn("a:headEnd"))
            if head_end is None:
                head_end = etree.SubElement(ln, qn("a:headEnd"))
            head_end.set("type", "triangle")
            head_end.set("w", "med")
            head_end.set("len", "med")

    @staticmethod
    def _anchor(node: "Node", side: str) -> Tuple[float, float]:
        if side == "top":
            return (node.cx, node.y)
        if side == "bottom":
            return (node.cx, node.bottom)
        if side == "left":
            return (node.x, node.cy)
        if side == "right":
            return (node.right, node.cy)
        if side == "topleft":
            return (node.x, node.y)
        if side == "topright":
            return (node.right, node.y)
        if side == "bottomleft":
            return (node.x, node.bottom)
        if side == "bottomright":
            return (node.right, node.bottom)
        raise ValueError(f"Unknown anchor side: {side}")

    def connect_line(
        self,
        src: str,
        dst: str,
        src_side: str = "right",
        dst_side: str = "left",
        kind: str = "straight",
        weight: float = 1.25,
        gap: float = 0.04,
    ) -> None:
        """Draw a thin arrow between two box edges.

        gap shortens both ends slightly so the head doesn't overlap the box border.
        """
        a = self.nodes[src]
        b = self.nodes[dst]
        x1, y1 = self._anchor(a, src_side)
        x2, y2 = self._anchor(b, dst_side)
        # Inset by gap along the direction src→dst so the arrowhead touches but doesn't overlap.
        dx, dy = x2 - x1, y2 - y1
        dist = max(1e-6, (dx * dx + dy * dy) ** 0.5)
        ux, uy = dx / dist, dy / dist
        x1 += ux * gap
        y1 += uy * gap
        x2 -= ux * gap
        y2 -= uy * gap
        self.line_arrow(x1, y1, x2, y2, weight=weight, kind=kind)

    # -- spine / branch helpers ---------------------------------------------
    def spine(
        self,
        items: List[Tuple[str, str, str]],
        x: float,
        y0: float,
        row_gap: float = 0.78,
        box_w: float = DEFAULT_BOX_W,
        box_h: float = DEFAULT_BOX_H,
        bold_ids: Optional[List[str]] = None,
    ) -> List[str]:
        """Stack a vertical spine of boxes connected by down-arrows.

        items is a list of (id, kind, label) tuples. Returns the list of ids.
        """
        bold_ids = set(bold_ids or [])
        ids: List[str] = []
        for i, (node_id, kind, label) in enumerate(items):
            self.box(
                node_id,
                kind,
                label,
                x=x,
                y=y0 + i * row_gap,
                w=box_w,
                h=box_h,
                bold=(node_id in bold_ids),
            )
            ids.append(node_id)
            if i > 0:
                self.connect_down(ids[i - 1], ids[i])
        return ids

    def branch_right(
        self,
        from_id: str,
        items: List[Tuple[str, str, str]],
        gap: float = 0.40,
        row_gap: float = 0.78,
        box_w: float = DEFAULT_BOX_W,
        box_h: float = DEFAULT_BOX_H,
    ) -> List[str]:
        """Place a vertical mini-chain to the right of from_id, joined by an arrow.

        First node aligns with from_id's row; subsequent nodes drop down.
        Returns the list of new ids.
        """
        anchor = self.nodes[from_id]
        x = anchor.right + gap
        ids: List[str] = []
        for i, (node_id, kind, label) in enumerate(items):
            self.box(node_id, kind, label, x=x, y=anchor.y + i * row_gap, w=box_w, h=box_h)
            ids.append(node_id)
            if i == 0:
                self.connect_right(from_id, node_id)
            else:
                self.connect_down(ids[i - 1], ids[i])
        return ids

    def branch_left(
        self,
        from_id: str,
        items: List[Tuple[str, str, str]],
        gap: float = 0.40,
        row_gap: float = 0.78,
        box_w: float = DEFAULT_BOX_W,
        box_h: float = DEFAULT_BOX_H,
    ) -> List[str]:
        anchor = self.nodes[from_id]
        x = anchor.x - gap - box_w
        ids: List[str] = []
        for i, (node_id, kind, label) in enumerate(items):
            self.box(node_id, kind, label, x=x, y=anchor.y + i * row_gap, w=box_w, h=box_h)
            ids.append(node_id)
            if i == 0:
                self.connect_left(from_id, node_id)
            else:
                self.connect_down(ids[i - 1], ids[i])
        return ids

    # -- text annotations ---------------------------------------------------
    def annotate(
        self,
        text: str,
        x: float,
        y: float,
        w: float = 1.7,
        h: float = 0.40,
        size: float = 9.5,
        font_name: str = "Cambria",
        color: RGBColor = COLOR_TEXT_LIGHT,
        align: PP_ALIGN = PP_ALIGN.LEFT,
    ) -> None:
        tb = self.slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.margin_left = Inches(0.03)
        tf.margin_right = Inches(0.03)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.color.rgb = color

    def _text(
        self,
        text: str,
        x: float,
        y: float,
        w: float,
        h: float,
        size: float = 11.0,
        bold: bool = False,
        font_name: str = "Calibri",
        color: RGBColor = COLOR_TEXT,
        align: PP_ALIGN = PP_ALIGN.LEFT,
    ) -> None:
        tb = self.slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.margin_left = Inches(0.04)
        tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.word_wrap = True
        for i, line in enumerate(text.split("\n")):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.alignment = align
            run = para.add_run()
            run.text = line
            run.font.name = font_name
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color

    # -- legend -------------------------------------------------------------
    def legend(self, x: float = 0.25, y: float = 4.95) -> None:
        """Compact legend in the lower-left under the Notes column."""
        self._text("Legend", x=x, y=y, w=1.4, h=0.25, size=12.0, bold=True)
        rows = [
            ("Custom scripts", COLOR_SCRIPT),
            ("CANlab Core functions", COLOR_FUNC),
            ("Key variables", COLOR_VAR),
            ("Files", COLOR_FILE),
        ]
        row_h = 0.28
        for i, (label, color) in enumerate(rows):
            top = y + 0.30 + i * row_h
            shape = self.slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(top), Inches(2.0), Inches(row_h - 0.04)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.color.rgb = COLOR_BORDER
            shape.line.width = Pt(0.5)
            tf = shape.text_frame
            tf.margin_left = Inches(0.06)
            tf.margin_right = Inches(0.04)
            tf.margin_top = Inches(0.02)
            tf.margin_bottom = Inches(0.02)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            run = para.add_run()
            run.text = label
            run.font.name = "Calibri"
            run.font.size = Pt(10.0)
            run.font.color.rgb = COLOR_TEXT

    # -- save ---------------------------------------------------------------
    def save(self, path: str) -> None:
        os.makedirs(os.path.dirname(path), exist_ok=True)
        self.prs.save(path)


# ---- Rendering -------------------------------------------------------------

def render_pptx_to_png(pptx_path: str, png_dir: str, crop: bool = True) -> str:
    """Convert a .pptx to .png using LibreOffice (soffice) and crop whitespace.

    Returns the path to the resulting cropped .png file.
    """
    os.makedirs(png_dir, exist_ok=True)
    soffice = shutil.which("soffice") or "/opt/homebrew/bin/soffice"
    with tempfile.TemporaryDirectory() as tmp:
        # soffice png export only emits the first slide as <name>.png
        subprocess.run(
            [soffice, "--headless", "--convert-to", "png", "--outdir", tmp, pptx_path],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        base = os.path.splitext(os.path.basename(pptx_path))[0]
        src_png = os.path.join(tmp, f"{base}.png")
        if not os.path.exists(src_png):
            raise RuntimeError(f"soffice did not produce {src_png}")
        out_png = os.path.join(png_dir, f"{base}.png")
        if crop:
            _crop_whitespace(src_png, out_png)
        else:
            shutil.copyfile(src_png, out_png)
    return out_png


def _crop_whitespace(src: str, dst: str, threshold: int = 245, pad: int = 12) -> None:
    """Crop near-white margins from an image, leaving a small padding."""
    from PIL import Image, ImageChops

    im = Image.open(src).convert("RGB")
    bg = Image.new("RGB", im.size, (255, 255, 255))
    diff = ImageChops.difference(im, bg)
    # Make the threshold tolerant of soffice's antialiased near-white pixels
    diff = diff.point(lambda p: 0 if p < (255 - threshold) else 255)
    bbox = diff.getbbox()
    if bbox is None:
        im.save(dst)
        return
    left = max(0, bbox[0] - pad)
    top = max(0, bbox[1] - pad)
    right = min(im.width, bbox[2] + pad)
    bottom = min(im.height, bbox[3] + pad)
    im.crop((left, top, right, bottom)).save(dst)
